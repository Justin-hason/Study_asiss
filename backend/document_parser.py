"""文档解析工具，支持 Word、PDF、PPT、Markdown、TXT 格式"""
import os
from typing import Optional, Dict, Any

# 文档类型到MIME类型的映射
DOCUMENT_TYPES = {
    ".docx": "application/msword",
    ".pdf": "application/pdf",
    ".pptx": "application/vnd.ms-powerpoint",
    ".doc": "application/msword",
    ".ppt": "application/vnd.ms-powerpoint",
    ".md": "text/markdown",
    ".txt": "text/plain",
}

TEXT_EXTENSIONS = {".md", ".txt"}


def get_file_type(filename: str) -> Optional[str]:
    """根据文件名获取文件类型"""
    ext = os.path.splitext(filename)[1].lower()
    return DOCUMENT_TYPES.get(ext)


def is_supported_document(filename: str) -> bool:
    """检查是否为支持的文档类型"""
    ext = os.path.splitext(filename)[1].lower()
    return ext in DOCUMENT_TYPES


def parse_document_content(file_path: str, mime_type: str) -> Dict[str, Any]:
    """
    解析文档内容

    Args:
        file_path: 文件路径
        mime_type: MIME类型

    Returns:
        包含文本内容和其他信息的字典
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    if ext in [".docx", ".doc"]:
        return parse_word(file_path)
    if ext in [".pptx", ".ppt"]:
        return parse_ppt(file_path)
    if ext in TEXT_EXTENSIONS:
        return parse_text(file_path)

    return {"error": "Unsupported file type", "text": ""}


def parse_pdf(file_path: str) -> Dict[str, Any]:
    """解析 PDF 文件"""
    try:
        import PyPDF2

        text_content = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    text_content.append({
                        "page": page_num + 1,
                        "text": text,
                    })

        full_text = "\n\n".join([page["text"] for page in text_content])

        return {
            "type": "pdf",
            "pages": len(text_content),
            "content": text_content,
            "text": full_text,
            "summary": f"PDF文档，共{len(text_content)}页",
        }
    except ImportError:
        return {"error": "PyPDF2 not installed", "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}


def parse_word(file_path: str) -> Dict[str, Any]:
    """解析 Word 文档"""
    try:
        from docx import Document

        doc = Document(file_path)

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        full_text = "\n".join(paragraphs)

        return {
            "type": "word",
            "paragraphs": len(paragraphs),
            "tables": len(tables),
            "content": {"paragraphs": paragraphs, "tables": tables},
            "text": full_text,
            "summary": f"Word文档，共{len(paragraphs)}个段落",
        }
    except ImportError:
        return {"error": "python-docx not installed", "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}


def parse_ppt(file_path: str) -> Dict[str, Any]:
    """解析 PowerPoint 文档"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)

        slides = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text

            slides.append({
                "slide_number": slide_num + 1,
                "title": title,
                "content": "\n".join(slide_text),
            })

        full_text = "\n\n".join([f"幻灯片 {slide['slide_number']}: {slide['content']}" for slide in slides])

        return {
            "type": "ppt",
            "slides": len(slides),
            "content": slides,
            "text": full_text,
            "summary": f"PowerPoint文档，共{len(slides)}张幻灯片",
        }
    except ImportError:
        return {"error": "python-pptx not installed", "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}


def parse_text(file_path: str) -> Dict[str, Any]:
    """解析 Markdown/TXT 文本文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="utf-8-sig") as f:
                text = f.read()
        except Exception as e:
            return {"error": str(e), "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}

    ext = os.path.splitext(file_path)[1].lower()
    lines = [line for line in text.splitlines() if line.strip()]

    return {
        "type": "markdown" if ext == ".md" else "text",
        "lines": len(lines),
        "content": {"lines": lines},
        "text": text,
        "summary": f"{'Markdown' if ext == '.md' else '文本'}文档，共{len(lines)}行",
    }
