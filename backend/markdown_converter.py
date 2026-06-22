"""文档转 Markdown 转换服务"""
import os
from typing import Optional, Dict, Any


def convert_to_markdown(file_path: str, mime_type: str) -> Optional[str]:
    """
    将文档转换为 Markdown 格式
    
    Args:
        file_path: 文件路径
        mime_type: MIME 类型
        
    Returns:
        Markdown 格式的文本内容，如果转换失败返回 None
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return convert_pdf_to_markdown(file_path)
    if ext in [".docx", ".doc"]:
        return convert_word_to_markdown(file_path)
    if ext in [".pptx", ".ppt"]:
        return convert_ppt_to_markdown(file_path)
    if ext in [".md", ".txt"]:
        return convert_text_to_markdown(file_path)
    
    return None


def convert_pdf_to_markdown(file_path: str) -> Optional[str]:
    """将 PDF 转换为 Markdown"""
    try:
        import PyPDF2
        
        markdown_parts = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    markdown_parts.append(f"## 第 {page_num} 页\n\n{text.strip()}\n")
        
        if not markdown_parts:
            return None
        
        return "\n".join(markdown_parts)
    except Exception as e:
        print(f"PDF 转换失败: {e}")
        return None


def convert_word_to_markdown(file_path: str) -> Optional[str]:
    """将 Word 文档转换为 Markdown"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        markdown_parts = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 根据样式判断标题级别
            if para.style.name.startswith("Heading"):
                level = int(para.style.name.replace("Heading", "").strip())
                level = min(level, 6)  # Markdown 最多支持 6 级标题
                markdown_parts.append(f"{'#' * level} {text}\n")
            else:
                # 普通段落
                markdown_parts.append(f"{text}\n")
        
        # 处理表格
        for table in doc.tables:
            markdown_parts.append("\n")
            # 表头
            header_row = table.rows[0]
            headers = [cell.text.strip() for cell in header_row.cells]
            markdown_parts.append("| " + " | ".join(headers) + " |\n")
            markdown_parts.append("| " + " | ".join(["---"] * len(headers)) + " |\n")
            
            # 数据行
            for row in table.rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                markdown_parts.append("| " + " | ".join(cells) + " |\n")
            markdown_parts.append("\n")
        
        if not markdown_parts:
            return None
        
        return "\n".join(markdown_parts)
    except Exception as e:
        print(f"Word 转换失败: {e}")
        return None


def convert_ppt_to_markdown(file_path: str) -> Optional[str]:
    """将 PowerPoint 转换为 Markdown"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            title = ""
            
            # 提取标题
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            
            slide_content.append(f"## 幻灯片 {slide_num}: {title or '无标题'}\n")
            
            # 提取内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # 避免重复标题
                        text = shape.text.strip()
                        slide_content.append(f"{text}\n")
            
            markdown_parts.append("\n".join(slide_content))
            markdown_parts.append("\n---\n")
        
        if not markdown_parts:
            return None
        
        return "\n".join(markdown_parts)
    except Exception as e:
        print(f"PPT 转换失败: {e}")
        return None


def convert_text_to_markdown(file_path: str) -> Optional[str]:
    """将文本文件转换为 Markdown（直接读取）"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="utf-8-sig") as f:
                content = f.read()
            return content
        except Exception as e:
            print(f"文本文件读取失败: {e}")
            return None
    except Exception as e:
        print(f"文本文件读取失败: {e}")
        return None


def outline_to_markdown(outline_content: Dict[str, Any]) -> str:
    """将大纲 JSON 转换为 Markdown"""
    if not outline_content or "sections" not in outline_content:
        return ""
    
    def convert_section(section: Dict[str, Any], level: int = 1) -> str:
        parts = []
        title = section.get("title", "")
        if title:
            parts.append(f"{'#' * min(level, 6)} {title}\n\n")
        
        children = section.get("children", [])
        for child in children:
            parts.append(convert_section(child, level + 1))
        
        return "".join(parts)
    
    sections = outline_content.get("sections", [])
    markdown_parts = []
    
    # 添加文档标题
    if "title" in outline_content:
        markdown_parts.append(f"# {outline_content['title']}\n\n")
    
    for section in sections:
        markdown_parts.append(convert_section(section, 2))
    
    return "\n".join(markdown_parts)


def note_to_markdown(note_content: str) -> str:
    """笔记内容已经是 Markdown 格式，直接返回"""
    return note_content


def knowledge_report_to_markdown(report_content: Dict[str, Any]) -> str:
    """将知识报告转换为 Markdown"""
    if not report_content:
        return ""
    
    markdown_parts = []
    
    # 标题
    if "title" in report_content:
        markdown_parts.append(f"# {report_content['title']}\n\n")
    
    # 摘要
    if "summary" in report_content:
        markdown_parts.append(f"## 摘要\n\n{report_content['summary']}\n\n")
    
    # 内容
    if "content" in report_content:
        content = report_content["content"]
        if isinstance(content, dict):
            # 递归转换内容
            def convert_content(obj: Any, level: int = 2) -> str:
                parts = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        if isinstance(value, (str, int, float)):
                            parts.append(f"{'#' * min(level, 6)} {key}\n\n{value}\n\n")
                        elif isinstance(value, list):
                            parts.append(f"{'#' * min(level, 6)} {key}\n\n")
                            for item in value:
                                if isinstance(item, dict):
                                    parts.append(convert_content(item, level + 1))
                                else:
                                    parts.append(f"- {item}\n")
                            parts.append("\n")
                        elif isinstance(value, dict):
                            parts.append(f"{'#' * min(level, 6)} {key}\n\n")
                            parts.append(convert_content(value, level + 1))
                return "".join(parts)
            
            markdown_parts.append(convert_content(content, 2))
        elif isinstance(content, str):
            markdown_parts.append(content)
    
    return "\n".join(markdown_parts)
